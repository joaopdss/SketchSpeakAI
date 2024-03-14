from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

class PowerPoint():

	def __init__(self):
		self.prs = Presentation()
		self.prs.slide_width = Inches(10)
		self.prs.slide_height = Inches(5.625)

	def create_pptx_slide(self, text, content_text, image_path, output_path):
		# Add slide
		slide_layout = self.prs.slide_layouts[1]
		slide = self.prs.slides.add_slide(slide_layout)

		# Set title instance
		title = slide.shapes.title

		# Set title position
		title.left = Inches(0.5)
		title.top = Inches(0.6)
		title.width = Inches(6)
		title.height = Inches(0.5)

		# Set content and font size
		title.text = text
		title.text_frame.paragraphs[0].font.size = Pt(22)

		# Set the alignment of the title text to left
		title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

		# Get content placeholder
		content_placeholder = slide.placeholders[1]

		# Set content positions
		content_placeholder.left = Inches(0.5)
		content_placeholder.top = Inches(1.5)
		content_placeholder.width = Inches(5)
		content_placeholder.height = Inches(3.5)

		# Put text of multiple lines
		for topic in content_text:
			p = content_placeholder.text_frame.add_paragraph()
			p.text = topic
			p.font.size = Pt(16)

		# Add the image to the slide
		img = slide.shapes.add_picture(image_path, Inches(6), Inches(2), width=Inches(3.3))

		# Calculate height to maintain aspect ratio
		aspect_ratio = img.image.size[1] / img.image.size[0]
		img.height = int(img.width * aspect_ratio)

		self.prs.save(output_path)

		return True


# Examples below
# title = "Lebron the GOAT of NBA"
# image_path = "../images/lebron.jpg"
# content_text = [
# 		"Lebron is considered for many the greatest basketball player of all time",
# 		"There people who compare him to Michael Jordan, but everyone knows who is the goat"
# 	]
# output_path = "teste.pptx"

# create_pptx_slide(title, content_text, image_path, output_path)