import openai
from record_audio import Audio
import time
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

### Just a bunch of text below that I dont wanna change today

# prompt = "Hey, give me some ideas for my new computer vision project!"


# response = get_response(prompt)

# print(response)

# response = get_response(True, "teste.wav")
# print(response)

def create_pptx_slide(text, content_text, image_path, output_path):
	prs = Presentation()

	# Set slide dimensions to widescreen (16:9)
	prs.slide_width = Inches(10)
	prs.slide_height = Inches(5.625)

	# Add slide
	slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(slide_layout)

	for shape in slide.shapes:
		if shape.shape_type == 17:  # Check if shape is a text box (shape_type 17 is for text boxes)
			slide.shapes._spTree.remove(shape._element)

	# Add the text to the slide
	title = slide.shapes.title
	title.text = text
	title.text_frame.paragraphs[0].font.size = Pt(20)

	max_text_width = int(prs.slide_width / 1.6)

	# Add the content container to the slide
	left = Inches(0.5)
	top = Inches(1.5)  # Adjust top position as needed
	width = max_text_width
	height = Inches(3.5)  # Adjust height as needed
	text_box = slide.shapes.add_textbox(left, top, width, height)

	# Add the content text to the container
	text_frame = text_box.text_frame
	text_frame.paragraphs[0].font.size = Pt(16)
	paragraph = text_frame.add_paragraph()
	paragraph.text = content_text

	# Enable wrapping
	text_frame.word_wrap = True

	# Add the image to the slide
	img = slide.shapes.add_picture(image_path, Inches(6.5), Inches(2), width=Inches(3))

	# Calculate height to maintain aspect ratio
	aspect_ratio = img.image.size[1] / img.image.size[0]
	img.height = int(img.width * aspect_ratio)

	prs.save(output_path)


title = "Lebron = GOAT"
image_path = "../images/lebron.jpg"
content_text = "Lebron is considered for many, the greatest basketball player of all time"
output_path = "teste.pptx"

create_pptx_slide(title, content_text, image_path, output_path)